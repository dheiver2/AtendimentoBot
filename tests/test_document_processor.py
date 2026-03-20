import os
import tempfile
import unittest

from docx import Document
from pptx import Presentation

import document_processor


class DocumentProcessorTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.original_pdfs_dir = document_processor.PDFS_DIR
        document_processor.PDFS_DIR = self.temp_dir.name

    def tearDown(self):
        document_processor.PDFS_DIR = self.original_pdfs_dir
        self.temp_dir.cleanup()

    def test_arquivo_suportado_reconhece_extensoes(self):
        self.assertTrue(document_processor.arquivo_suportado("manual.pdf"))
        self.assertTrue(document_processor.arquivo_suportado("roteiro.DOCX"))
        self.assertFalse(document_processor.arquivo_suportado("script.exe"))

    def test_processa_arquivo_texto_e_salva_no_disco(self):
        chunks = document_processor.processar_documento(
            10,
            "base.txt",
            b"linha 1\nlinha 2\nlinha 3",
        )

        caminho = os.path.join(self.temp_dir.name, "10", "base.txt")

        self.assertTrue(os.path.exists(caminho))
        self.assertGreaterEqual(len(chunks), 1)
        self.assertIn("linha 1", chunks[0])

    def test_processa_docx_salvo(self):
        caminho = os.path.join(self.temp_dir.name, "base.docx")
        doc = Document()
        doc.add_paragraph("Politica de troca em ate 7 dias.")
        doc.save(caminho)

        chunks = document_processor.processar_documento_salvo(caminho)

        self.assertGreaterEqual(len(chunks), 1)
        self.assertIn("Politica de troca", chunks[0])

    def test_processa_pptx_salvo(self):
        caminho = os.path.join(self.temp_dir.name, "apresentacao.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Atendimento"
        slide.placeholders[1].text = "Horario de segunda a sexta."
        prs.save(caminho)

        chunks = document_processor.processar_documento_salvo(caminho)

        self.assertGreaterEqual(len(chunks), 1)
        self.assertIn("Horario de segunda a sexta", chunks[0])

    def test_rejeita_formato_nao_suportado(self):
        with self.assertRaises(ValueError):
            document_processor.processar_documento(1, "arquivo.zip", b"conteudo")
