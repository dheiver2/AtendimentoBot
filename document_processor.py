import os
from pathlib import Path

from docx import Document as WordDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader

from config import PDFS_DIR
from validators import sanitizar_nome_arquivo, validar_tamanho_documento

_DEFAULT_CHUNK_SIZE = 700
_DEFAULT_CHUNK_OVERLAP = 80

SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".pptx": "PowerPoint",
    ".txt": "texto",
    ".md": "Markdown",
    ".csv": "CSV",
}


def listar_formatos_suportados() -> str:
    """Retorna a lista curta dos formatos suportados."""
    return ", ".join(SUPPORTED_EXTENSIONS)


def arquivo_suportado(nome_arquivo: str) -> bool:
    """Verifica se a extensão do arquivo é suportada."""
    return Path(nome_arquivo).suffix.lower() in SUPPORTED_EXTENSIONS


def salvar_documento(empresa_id: int, nome_arquivo: str, conteudo_bytes: bytes) -> str:
    """Salva o documento no disco e retorna o caminho."""
    pasta_empresa = os.path.join(PDFS_DIR, str(empresa_id))
    os.makedirs(pasta_empresa, exist_ok=True)
    caminho = os.path.join(pasta_empresa, nome_arquivo)
    with open(caminho, "wb") as f:
        f.write(conteudo_bytes)
    return caminho


def _extrair_texto_pdf(caminho_arquivo: str) -> str:
    reader = PdfReader(caminho_arquivo)
    partes = []
    for pagina in reader.pages:
        texto_pagina = pagina.extract_text()
        if texto_pagina:
            partes.append(texto_pagina.strip())
    return "\n\n".join(partes)


def _extrair_texto_docx(caminho_arquivo: str) -> str:
    doc = WordDocument(caminho_arquivo)
    partes = [paragrafo.text.strip() for paragrafo in doc.paragraphs if paragrafo.text.strip()]

    for tabela in doc.tables:
        for linha in tabela.rows:
            celulas = [celula.text.strip() for celula in linha.cells if celula.text.strip()]
            if celulas:
                partes.append(" | ".join(celulas))

    return "\n".join(partes)


def _extrair_texto_pptx(caminho_arquivo: str) -> str:
    apresentacao = Presentation(caminho_arquivo)
    partes = []

    for slide in apresentacao.slides:
        for shape in slide.shapes:
            texto = getattr(shape, "text", "")
            if texto and texto.strip():
                partes.append(texto.strip())

    return "\n\n".join(partes)


def _extrair_texto_simples(caminho_arquivo: str) -> str:
    for encoding in ["utf-8-sig", "utf-8", "latin-1"]:
        try:
            with open(caminho_arquivo, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    raise ValueError("Não foi possível ler o arquivo de texto com uma codificação suportada.")


def extrair_texto_documento(caminho_arquivo: str, extensao: str) -> str:
    """Extrai texto de um documento suportado."""
    if extensao == ".pdf":
        return _extrair_texto_pdf(caminho_arquivo)
    if extensao == ".docx":
        return _extrair_texto_docx(caminho_arquivo)
    if extensao == ".pptx":
        return _extrair_texto_pptx(caminho_arquivo)
    if extensao in {".txt", ".md", ".csv"}:
        return _extrair_texto_simples(caminho_arquivo)

    raise ValueError(
        f"Formato não suportado. Envie um dos seguintes formatos: {listar_formatos_suportados()}."
    )


def dividir_texto_em_chunks(
    texto: str,
    chunk_size: int = _DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = _DEFAULT_CHUNK_OVERLAP,
) -> list[str]:
    """Divide o texto em pedaços menores para indexação."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(texto)


def processar_documento(empresa_id: int, nome_arquivo: str, conteudo_bytes: bytes) -> list[str]:
    """Pipeline completo: valida, salva, extrai texto e divide em chunks."""
    nome_arquivo = sanitizar_nome_arquivo(nome_arquivo)
    validar_tamanho_documento(len(conteudo_bytes), nome_arquivo)

    extensao = Path(nome_arquivo).suffix.lower()
    if extensao not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Formato não suportado. Envie um dos seguintes formatos: {listar_formatos_suportados()}."
        )

    caminho = salvar_documento(empresa_id, nome_arquivo, conteudo_bytes)
    texto = extrair_texto_documento(caminho, extensao)
    if not texto.strip():
        raise ValueError(
            "Não foi possível extrair texto do arquivo. Verifique se ele contém texto selecionável."
        )

    return dividir_texto_em_chunks(texto)


def processar_documento_salvo(caminho_arquivo: str) -> list[str]:
    """Processa um documento já salvo no disco e retorna os chunks."""
    extensao = Path(caminho_arquivo).suffix.lower()
    if extensao not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Formato não suportado. Envie um dos seguintes formatos: {listar_formatos_suportados()}."
        )

    texto = extrair_texto_documento(caminho_arquivo, extensao)
    if not texto.strip():
        raise ValueError(
            "Não foi possível extrair texto do arquivo. Verifique se ele contém texto selecionável."
        )

    return dividir_texto_em_chunks(texto)
